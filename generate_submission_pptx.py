"""
Build a single-slide PowerPoint cover page (dashboard-style) for assignment submission.

Usage:
  1. Set STREAMLIT_APP_URL below to your public Streamlit Cloud link.
  2. Run: python generate_submission_pptx.py
  3. Open submission_dashboard_cover.pptx in PowerPoint on your laptop.

Requires: pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# EDIT THIS — your deployed app URL (must be reachable in a browser)
# ---------------------------------------------------------------------------
STREAMLIT_APP_URL = "https://YOUR_APP_NAME.streamlit.app"

OUT_PATH = Path(__file__).resolve().parent / "submission_dashboard_cover.pptx"

# Dashboard dark theme (matches app defaults)
BG = RGBColor(0x0C, 0x11, 0x18)
SIDEBAR = RGBColor(0x0F, 0x14, 0x19)
FG = RGBColor(0xE5, 0xE7, 0xEB)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
H1 = RGBColor(0xF0, 0xF9, 0xFF)
H2 = RGBColor(0xBA, 0xE6, 0xFD)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
RIBBON = RGBColor(0x0E, 0xA5, 0xE9)


def _set_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape) -> None:
    shape.line.fill.background()


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)

    # Full slide background
    bg = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )  # 1 = MSO_SHAPE.RECTANGLE
    _set_fill(bg, BG)
    _no_line(bg)
    # Send to back (z-order): first shape is often back - good

    # Ribbon bar
    ribbon_h = Inches(0.06)
    rib = slide.shapes.add_shape(1, 0, 0, prs.slide_width, ribbon_h)
    _set_fill(rib, RIBBON)
    _no_line(rib)

    # Title block
    left = Inches(0.55)
    top = Inches(0.35)
    tw = prs.slide_width - Inches(1.1)
    title_box = slide.shapes.add_textbox(left, top, tw, Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NASA C-MAPSS Turbofan Engine — PHM Dashboard"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = H1
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = "Assignment 2 · Interactive fleet health & RUL briefing (Streamlit)"
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(6)
    p2.font.name = "Calibri"

    # Mock layout: sidebar + main
    y0 = Inches(1.45)
    side_w = Inches(2.35)
    gap = Inches(0.08)
    main_left = left + side_w + gap
    main_w = prs.slide_width - main_left - Inches(0.55)
    h = Inches(2.55)

    side = slide.shapes.add_shape(1, left, y0, side_w, h)
    _set_fill(side, SIDEBAR)
    side.line.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
    side.line.width = Pt(0.5)

    stf = slide.shapes.add_textbox(left + Inches(0.12), y0 + Inches(0.12), side_w - Inches(0.2), h - Inches(0.2))
    st = stf.text_frame
    st.word_wrap = True
    sp = st.paragraphs[0]
    sp.text = "APPEARANCE"
    sp.font.size = Pt(8)
    sp.font.bold = True
    sp.font.color.rgb = MUTED
    sp = st.add_paragraph()
    sp.text = "Dark · Light"
    sp.font.size = Pt(10)
    sp.font.color.rgb = FG
    sp.space_before = Pt(4)
    sp = st.add_paragraph()
    sp.text = "DATA SOURCE"
    sp.font.size = Pt(8)
    sp.font.bold = True
    sp.font.color.rgb = MUTED
    sp.space_before = Pt(10)
    sp = st.add_paragraph()
    sp.text = "CMAPSSData · FD001–FD004 · Fleet test/train"
    sp.font.size = Pt(9)
    sp.font.color.rgb = FG
    sp.space_before = Pt(4)

    main = slide.shapes.add_textbox(main_left, y0, main_w, h)
    mf = main.text_frame
    mf.word_wrap = True
    mp = mf.paragraphs[0]
    mp.text = "Home — three engineering pillars"
    mp.font.size = Pt(14)
    mp.font.bold = True
    mp.font.color.rgb = H2
    mp.font.name = "Calibri"

    bullets = (
        "Physical system insight — turbofan context, 26 channels per cycle, degradation in time series.\n"
        "Model trade-offs — CNN / Transformer vs PINN and hybrid stacks.\n"
        "Compute & data reality — data quality, drift, deployment cost (GEMM, KV cache, quantization)."
    )
    mp2 = mf.add_paragraph()
    mp2.text = bullets
    mp2.font.size = Pt(11)
    mp2.font.color.rgb = FG
    mp2.space_before = Pt(8)
    mp2.line_spacing = 1.15
    mp2.font.name = "Calibri"

    # Tab strip (mock)
    tab_y = y0 + h + Inches(0.08)
    tabs = (
        "Fleet overview",
        "Exploratory analysis",
        "Data audit",
        "Architecture",
        "DL lab",
        "Compute",
    )
    x = main_left
    for i, t in enumerate(tabs):
        twd = Inches(1.45) if i < 5 else Inches(1.2)
        tb = slide.shapes.add_shape(1, x, tab_y, twd, Inches(0.28))
        _set_fill(tb, RGBColor(0x0F, 0x17, 0x2A))
        tb.line.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
        tb.line.width = Pt(0.25)
        ttf = tb.text_frame
        ttf.paragraphs[0].text = t
        ttf.paragraphs[0].font.size = Pt(7)
        ttf.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8) if i else H1
        ttf.paragraphs[0].font.bold = True if i == 0 else False
        ttf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += twd + Inches(0.06)

    # Group table + link
    ty = tab_y + Inches(0.45)
    tbl = slide.shapes.add_table(4, 2, left, ty, Inches(6.2), Inches(1.05)).table
    hdr = ("Name", "Student ID")
    for c, h in enumerate(hdr):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = H2
    rows = [
        ("Kodie Amo Kwame", "LS2525226"),
        ("Sumara Alfred Salifu", "LS2525245"),
        ("Peta Mimi Precious", "LS2525255"),
    ]
    for r, (name, sid) in enumerate(rows, start=1):
        tbl.cell(r, 0).text = name
        tbl.cell(r, 1).text = sid
        for c in range(2):
            for p in tbl.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = FG

    # Style table cells background (optional - simple)
    for r in range(4):
        for c in range(2):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    # Live URL
    link_top = ty + Inches(1.18)
    lx = left
    lw = prs.slide_width - Inches(1.1)
    lbox = slide.shapes.add_textbox(lx, link_top, lw, Inches(0.95))
    lf = lbox.text_frame
    lf.word_wrap = True
    lp = lf.paragraphs[0]
    lp.text = "Public Streamlit app (for lecturer): "
    lp.font.size = Pt(11)
    lp.font.color.rgb = MUTED
    lp.font.name = "Calibri"

    lp2 = lf.add_paragraph()
    run = lp2.add_run()
    run.text = STREAMLIT_APP_URL
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Calibri"
    run.hyperlink.address = STREAMLIT_APP_URL

    lp3 = lf.add_paragraph()
    lp3.text = (
        "Beihang University · RCSSTEAP, China — Replace YOUR_APP_NAME in generate_submission_pptx.py "
        "if the link above is still a placeholder."
    )
    lp3.font.size = Pt(9)
    lp3.font.color.rgb = MUTED
    lp3.space_before = Pt(6)

    prs.save(out_path := OUT_PATH)
    print(f"Wrote: {out_path}")


if __name__ == "__main__":
    main()
